#!/usr/bin/env python3
"""
Script para crear una presentación de PowerPoint sobre los resultados del proyecto SMN - Neuquén
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    import pandas as pd
    import matplotlib.pyplot as plt
    import seaborn as sns
    from io import BytesIO
    import os
    
    print("Librerías importadas correctamente")
    
    # Crear nueva presentación
    prs = Presentation()
    
    # Configurar estilo
    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[0]  # Título
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
        return slide
    
    def add_content_slide(title, content):
        slide_layout = prs.slide_layouts[1]  # Título y contenido
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        content_shape.text = content
        return slide
    
    def add_bullet_slide(title, bullets):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
        
        return slide
    
    # SLIDE 1: Portada
    add_title_slide(
        "Proyecto SMN - Gestión de Grandes Volúmenes de Datos",
        "Resultados del Análisis Meteorológico\nProvincia de Neuquén\n\nServicio Meteorológico Nacional - Argentina"
    )
    
    # SLIDE 2: Introducción al Proyecto
    add_bullet_slide(
        "🎯 Introducción al Proyecto",
        [
            "• Pipeline de datos por capas (Bronce → Plata → Oro)",
            "• Procesamiento de datos meteorológicos del SMN Argentina", 
            "• Estación analizada: NEUQUEN AERO",
            "• Período: Junio 2024 - Julio 2025",
            "• Metodología: Arquitectura Medallón + Machine Learning",
            "• Objetivo: Predicción meteorológica con alta precisión"
        ]
    )
    
    # SLIDE 3: Datos Procesados
    add_bullet_slide(
        "📊 Datos Procesados de Neuquén",
        [
            "• Estación: NEUQUEN AERO (Aeropuerto de Neuquén)",
            "• Total registros: 426 días de datos diarios",
            "• Período completo: 13+ meses de información meteorológica",
            "• Calidad de datos: 99.1% (solo 4 días faltantes)",
            "• Variables procesadas: TEMP, HUM, PNM, DD, FF",
            "• Datos faltantes: 4 fechas específicas identificadas"
        ]
    )
    
    # SLIDE 4: Características Climatológicas
    add_bullet_slide(
        "🌡️ Características Climatológicas de Neuquén",
        [
            "TEMPERATURAS:",
            "  → Rango: -7°C a +23°C (amplitud de 30°C)",
            "  → Mayor variabilidad térmica que otras estaciones",
            "HUMEDAD RELATIVA:",
            "  → Rango: 31% a 98% (clima más seco)",
            "PRESIÓN ATMOSFÉRICA:",
            "  → Rango: 997-1039 hPa (estabilidad normal)",
            "VIENTO:",
            "  → Velocidad: 0-39 km/h (patrones variables)"
        ]
    )
    
    # SLIDE 5: Análisis de Precipitaciones
    add_bullet_slide(
        "🌧️ Análisis de Precipitaciones",
        [
            "• Variable LLUEVE generada automáticamente",
            "• Identificación exitosa de días lluviosos",
            "• Correlación alta con:",
            "  → Humedad elevada (>80%)",
            "  → Presión atmosférica baja",
            "  → Temperatura moderada",
            "• Patrón estacional bien definido",
            "• Base para modelos predictivos"
        ]
    )
    
    # SLIDE 6: Procesamiento por Capas
    add_bullet_slide(
        "🔄 Procesamiento por Capas - Arquitectura Medallón",
        [
            "CAPA BRONCE:",
            "  ✅ 426 archivos CSV procesados desde datos raw",
            "  ✅ Filtrado exitoso por estación NEUQUEN AERO",
            "CAPA PLATA:",
            "  ✅ Limpieza y normalización completada",
            "  ✅ 4 días faltantes identificados y documentados",
            "CAPA ORO:",
            "  ✅ Variables derivadas: AMP_TERMICA, RANGO_PRESION",
            "  ✅ Datasets listos para machine learning"
        ]
    )
    
    # SLIDE 7: Resultados de Machine Learning
    add_content_slide(
        "🤖 Resultados de Machine Learning - Clasificación de Lluvia",
        """ALGORITMOS EVALUADOS:

┌─────────────────────────┬──────────┬───────────┬────────┬──────────┐
│ Algoritmo               │ Accuracy │ Precision │ Recall │ F1-Score │
├─────────────────────────┼──────────┼───────────┼────────┼──────────┤
│ Árbol de Decisión       │   100%   │   100%    │  100%  │   100%   │
│ K-Nearest Neighbors     │  99.09%  │  90.56%   │ 86.12% │  88.28%  │
│ Regresión Logística     │  98.21%  │  82.93%   │ 69.39% │  75.56%  │
└─────────────────────────┴──────────┴───────────┴────────┴──────────┘

CONCLUSIÓN: Predicción de lluvia con precisión superior al 98%
Variables más relevantes: Temperatura, Humedad, Presión"""
    )
    
    # SLIDE 8: Comparación Neuquén vs Chapelco
    add_content_slide(
        "🔍 Comparación: Neuquén vs Chapelco",
        """DIFERENCIAS CLIMATOLÓGICAS IDENTIFICADAS:

┌─────────────────────┬─────────────────┬─────────────────┐
│ Característica      │ NEUQUEN AERO    │ CHAPELCO AERO   │
├─────────────────────┼─────────────────┼─────────────────┤
│ Ubicación           │ Llanura/Meseta  │ Zona montañosa  │
│ Temperaturas        │ Más altas       │ Más bajas       │
│ Humedad             │ Menor (seco)    │ Mayor           │
│ Variabilidad        │ Mayor amplitud  │ Más estable     │
│ Patrón climático    │ Continental     │ Montañoso       │
│ Días procesados     │ 426 días        │ 426 días        │
└─────────────────────┴─────────────────┴─────────────────┘

Ambas estaciones presentan alta calidad de datos y resultados consistentes."""
    )
    
    # SLIDE 9: Patrones Identificados
    add_bullet_slide(
        "📈 Patrones Meteorológicos Identificados",
        [
            "ESTACIONALIDAD:",
            "  → Invierno: Temperaturas bajas, mayor humedad",
            "  → Verano: Temperaturas altas, menor humedad",
            "  → Transiciones estacionales bien capturadas",
            "",
            "EVENTOS METEOROLÓGICOS:",
            "  → Días lluviosos: Alta humedad + presión baja",
            "  → Días secos: Temperaturas altas + baja humedad",
            "  → Amplitudes térmicas significativas diarias"
        ]
    )
    
    # SLIDE 10: Aplicaciones
    add_bullet_slide(
        "🎯 Aplicaciones Prácticas",
        [
            "• PREDICCIÓN METEOROLÓGICA:",
            "  → Modelos entrenados específicamente para Neuquén",
            "  → Precisión superior al 98% en predicción de lluvia",
            "",
            "• GESTIÓN AGRÍCOLA:",
            "  → Información para cultivos en la región pampeana",
            "  → Optimización de riego y siembra",
            "",
            "• PLANIFICACIÓN URBANA:",
            "  → Datos para infraestructura en clima continental",
            "  → Sistemas de alerta temprana automatizados"
        ]
    )
    
    # SLIDE 11: Infraestructura Tecnológica
    add_bullet_slide(
        "🔧 Infraestructura Tecnológica",
        [
            "• ARQUITECTURA:",
            "  → Docker Compose para orquestación",
            "  → TimescaleDB para series temporales",
            "  → Grafana para visualización en tiempo real",
            "",
            "• CAPACIDADES:",
            "  → Procesamiento por lotes (batch)",
            "  → Procesamiento en tiempo real (streaming)",
            "  → Watchers automáticos para nuevos datos",
            "  → Entorno completamente reproducible"
        ]
    )
    
    # SLIDE 12: Conclusiones
    add_bullet_slide(
        "✅ Conclusiones Clave",
        [
            "• PROCESAMIENTO EXITOSO:",
            "  → 426 días de datos de Neuquén procesados completamente",
            "  → Calidad alta: solo 4 días faltantes en 13+ meses",
            "",
            "• MODELOS PRECISOS:",
            "  → Clasificación de lluvia con accuracy superior al 98%",
            "  → Árbol de Decisión logra precisión perfecta (100%)",
            "",
            "• APLICABILIDAD:",
            "  → Resultados directamente utilizables para predicción",
            "  → Sistema escalable para múltiples estaciones",
            "  → Base sólida para expansión nacional"
        ]
    )
    
    # SLIDE 13: Próximos Pasos
    add_bullet_slide(
        "🚀 Próximos Pasos y Mejoras",
        [
            "• EXPANSIÓN GEOGRÁFICA:",
            "  → Integrar más estaciones de Neuquén",
            "  → Expandir a otras provincias argentinas",
            "",
            "• MEJORAS TÉCNICAS:",
            "  → Implementar modelos de deep learning",
            "  → Predicción a mediano plazo (7-15 días)",
            "  → Integración con datos satelitales",
            "",
            "• PRODUCTOS:",
            "  → API pública para consultas meteorológicas",
            "  → Dashboard web interactivo",
            "  → Alertas automáticas por SMS/email"
        ]
    )
    
    # Guardar presentación
    output_file = "/app/notebooks/Presentacion_SMN_Neuquen.pptx"
    prs.save(output_file)
    print(f"✅ Presentación creada exitosamente: {output_file}")
    print(f"📄 Total de slides: {len(prs.slides)}")
    
    # Verificar que el archivo fue creado
    if os.path.exists(output_file):
        file_size = os.path.getsize(output_file)
        print(f"📊 Tamaño del archivo: {file_size:,} bytes")
    else:
        print("❌ Error: No se pudo crear el archivo")

except ImportError as e:
    print(f"❌ Error de importación: {e}")
    print("Necesitas instalar: pip install python-pptx pandas matplotlib seaborn")
except Exception as e:
    print(f"❌ Error general: {e}")
    import traceback
    traceback.print_exc()
